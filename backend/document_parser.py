"""
Document parsing service.

This module is intentionally decoupled from agent logic so it can be
plugged into RAG evaluation frameworks (RAGas, Arize Phoenix, etc.)
to measure retrieval context precision independently.

Supported formats (no external heavy dependencies):
  - PDF        → pypdf
  - DOCX       → python-docx
  - XLSX/XLS   → pandas + openpyxl
  - CSV        → pandas
  - PPTX/PPT   → python-pptx
  - Plain text / Markdown / JSON / HTML → direct decode
  - Google Workspace exports → plain text
"""

import io
import logging
import os
from dataclasses import dataclass, field
from typing import Optional

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Data container
# ---------------------------------------------------------------------------

@dataclass
class ParsedDocument:
    """Structured result from the parsing pipeline."""
    source_id: str                  # File ID or URI
    file_name: str                  # Original file name
    mime_type: str                  # MIME type of the source file
    chunks: list[str] = field(default_factory=list)  # Text chunks
    metadata: dict = field(default_factory=dict)      # Arbitrary metadata
    error: Optional[str] = None     # Set when parsing fails

    @property
    def full_text(self) -> str:
        """Join chunks into a single string for simple use-cases."""
        return "\n\n".join(self.chunks)

    @property
    def is_empty(self) -> bool:
        return not any(c.strip() for c in self.chunks)


# ---------------------------------------------------------------------------
# MIME-type helpers
# ---------------------------------------------------------------------------

PDF_MIMES = {
    "application/pdf",
}

WORD_MIMES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # .docx
    "application/msword",  # .doc
}

SPREADSHEET_MIMES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",  # .xlsx
    "application/vnd.ms-excel",  # .xls
    "text/csv",
    "application/csv",
}

PRESENTATION_MIMES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # .pptx
    "application/vnd.ms-powerpoint",  # .ppt
}

PLAIN_TEXT_MIMES = {
    "text/plain",
    "text/markdown",
    "text/x-markdown",
    "text/x-rst",
    "application/json",
    "text/html",
    "application/xml",
    "text/xml",
}

GOOGLE_DOC_MIMES = {
    "application/vnd.google-apps.document",
    "application/vnd.google-apps.spreadsheet",
    "application/vnd.google-apps.presentation",
    "application/vnd.google-apps.drawing",
}

CHUNK_SIZE = 3000      # Characters per chunk
CHUNK_OVERLAP = 200    # Overlap between chunks


def _chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split text into overlapping chunks for downstream retrieval."""
    if not text.strip():
        return []
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        if end == len(text):
            break
        start = end - overlap
    return chunks


# ---------------------------------------------------------------------------
# Format-specific parsers
# ---------------------------------------------------------------------------

def _parse_pdf(raw_bytes: bytes, file_name: str) -> list[str]:
    """Extract text from a PDF using pypdf."""
    try:
        import pypdf  # type: ignore

        reader = pypdf.PdfReader(io.BytesIO(raw_bytes))
        pages_text: list[str] = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages_text.append(f"[Page {i + 1}]\n{text.strip()}")

        full_text = "\n\n".join(pages_text)
        logger.info("PDF '%s': extracted %d pages, %d chars", file_name, len(reader.pages), len(full_text))
        return _chunk_text(full_text)

    except Exception as exc:
        logger.error("pypdf failed for '%s': %s", file_name, exc)
        return []


def _parse_docx(raw_bytes: bytes, file_name: str) -> list[str]:
    """Extract text from a .docx file using python-docx."""
    try:
        import docx  # type: ignore

        doc = docx.Document(io.BytesIO(raw_bytes))
        parts: list[str] = []

        # Paragraphs
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text.strip())

        # Tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    parts.append(row_text)

        full_text = "\n\n".join(parts)
        logger.info("DOCX '%s': extracted %d paragraphs/rows", file_name, len(parts))
        return _chunk_text(full_text)

    except Exception as exc:
        logger.error("python-docx failed for '%s': %s", file_name, exc)
        return []


def _parse_spreadsheet(raw_bytes: bytes, file_name: str, mime_type: str) -> list[str]:
    """Extract text from CSV/Excel spreadsheets using pandas."""
    try:
        import pandas as pd  # type: ignore

        chunks: list[str] = []

        if mime_type in {"text/csv", "application/csv"} or file_name.lower().endswith(".csv"):
            df = pd.read_csv(io.BytesIO(raw_bytes))
            chunks.append(f"CSV: {file_name}\n\n{df.to_string(index=False)}")
        else:
            # Excel – iterate sheets
            xl = pd.ExcelFile(io.BytesIO(raw_bytes))
            for sheet_name in xl.sheet_names:
                df = xl.parse(sheet_name)
                text = f"Sheet: {sheet_name}\n\n{df.to_string(index=False)}"
                chunks.append(text)

        all_text = "\n\n---\n\n".join(chunks)
        logger.info("Spreadsheet '%s': extracted %d chars", file_name, len(all_text))
        return _chunk_text(all_text)

    except Exception as exc:
        logger.error("pandas failed for '%s': %s", file_name, exc)
        return []


def _parse_pptx(raw_bytes: bytes, file_name: str) -> list[str]:
    """Extract text from a .pptx file using python-pptx."""
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation(io.BytesIO(raw_bytes))
        slides_text: list[str] = []

        for i, slide in enumerate(prs.slides):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if text:
                            slide_parts.append(text)
            if slide_parts:
                slides_text.append(f"[Slide {i + 1}]\n" + "\n".join(slide_parts))

        full_text = "\n\n".join(slides_text)
        logger.info("PPTX '%s': extracted %d slides, %d chars", file_name, len(prs.slides), len(full_text))
        return _chunk_text(full_text)

    except Exception as exc:
        logger.error("python-pptx failed for '%s': %s", file_name, exc)
        return []


def _parse_plain_text(raw_bytes: bytes, file_name: str) -> list[str]:
    """Decode and chunk plain text / markdown / JSON files."""
    try:
        text = raw_bytes.decode("utf-8", errors="replace")
        logger.info("Plain text '%s': %d chars", file_name, len(text))
        return _chunk_text(text)
    except Exception as exc:
        logger.error("Plain text decode failed for '%s': %s", file_name, exc)
        return []


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def parse_document(
    raw_bytes: bytes,
    file_name: str,
    mime_type: str,
    source_id: str = "",
) -> ParsedDocument:
    """
    Parse raw file bytes into a :class:`ParsedDocument`.

    This function is the single entry point for the parsing pipeline.
    It dispatches to a format-specific parser based on the MIME type.

    Parameters
    ----------
    raw_bytes:
        The raw file content as bytes.
    file_name:
        Original file name (used for logging and extension-based fallback).
    mime_type:
        MIME type string (e.g. ``"application/pdf"``).
    source_id:
        Identifier for the source (e.g. Google Drive file ID).

    Returns
    -------
    ParsedDocument
        Structured container with text chunks and metadata.
    """
    result = ParsedDocument(
        source_id=source_id,
        file_name=file_name,
        mime_type=mime_type,
        metadata={"source_id": source_id, "file_name": file_name, "mime_type": mime_type},
    )

    try:
        fname_lower = file_name.lower()

        # --- PDF ---
        if mime_type in PDF_MIMES or fname_lower.endswith(".pdf"):
            result.chunks = _parse_pdf(raw_bytes, file_name)

        # --- Word ---
        elif mime_type in WORD_MIMES or fname_lower.endswith((".docx", ".doc")):
            result.chunks = _parse_docx(raw_bytes, file_name)

        # --- Spreadsheets / CSV ---
        elif mime_type in SPREADSHEET_MIMES or fname_lower.endswith((".xlsx", ".xls", ".csv")):
            result.chunks = _parse_spreadsheet(raw_bytes, file_name, mime_type)

        # --- Presentations ---
        elif mime_type in PRESENTATION_MIMES or fname_lower.endswith((".pptx", ".ppt")):
            result.chunks = _parse_pptx(raw_bytes, file_name)

        # --- Google Workspace exports (already decoded to text by MCP) ---
        elif mime_type in GOOGLE_DOC_MIMES:
            result.chunks = _parse_plain_text(raw_bytes, file_name)

        # --- Plain text / markdown / JSON / HTML ---
        elif mime_type in PLAIN_TEXT_MIMES:
            result.chunks = _parse_plain_text(raw_bytes, file_name)

        # --- Try plain text as last resort ---
        else:
            logger.warning(
                "Unknown MIME type '%s' for '%s' – attempting plain text decode",
                mime_type, file_name,
            )
            result.chunks = _parse_plain_text(raw_bytes, file_name)

    except Exception as exc:
        logger.exception("Unexpected error parsing '%s'", file_name)
        result.error = str(exc)

    if result.is_empty and not result.error:
        result.error = f"No text could be extracted from '{file_name}' (type: {mime_type})."

    return result
